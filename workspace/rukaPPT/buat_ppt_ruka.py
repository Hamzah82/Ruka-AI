"""
🐢 Ruka AI — Pembuat PPT Pengenalan Diri
Menghasilkan file Ruka_AI_Pengenalan.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Warna Tema ───
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # Biru gelap navy
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)    # Cyan terang
ACCENT_PURPLE = RGBColor(0xA8, 0x5C, 0xFF)  # Ungu soft
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)   # Hijau mint
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)  # Oranye soft
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def dark_bg(slide):
    """Set background gelap untuk slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Tambah shape dengan warna"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Tambah text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_circle(slide, left, top, size, color):
    """Tambah lingkaran"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    """Tambah bullet list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_icon_text(slide, left, top, icon, text, font_size=14, color=WHITE):
    """Tambah teks dengan icon emoji di depannya"""
    full_text = f"{icon}  {text}"
    txBox = slide.shapes.add_textbox(left, top, Inches(5), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = full_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return txBox


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — JUDUL / COVER
# ═══════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
dark_bg(slide1)

# Dekorasi lingkaran besar di kiri
add_circle(slide1, Inches(-1), Inches(-1), Inches(4), RGBColor(0x16, 0x21, 0x37))
add_circle(slide1, Inches(9), Inches(4), Inches(5), RGBColor(0x16, 0x21, 0x37))

# Garis aksen di kiri
accent_line = add_shape(slide1, Inches(1.2), Inches(1.5), Inches(0.08), Inches(4.5), ACCENT_BLUE, MSO_SHAPE.RECTANGLE)

# Judul utama
add_text_box(slide1, Inches(1.2), Inches(1.8), Inches(10), Inches(1.2),
             "🐢 RUKA AI", font_size=60, color=WHITE, bold=True)

# Sub judul
add_text_box(slide1, Inches(1.2), Inches(3.0), Inches(10), Inches(0.8),
             "AI Agent Berbentuk Kura-Kura", font_size=32, color=ACCENT_BLUE, bold=True)

# Deskripsi singkat
add_text_box(slide1, Inches(1.2), Inches(3.9), Inches(10), Inches(1),
             "Asisten AI berbasis terminal yang bijaksana, sabar, dan teliti.\nBerjalan di mesin lokal — private, aman, dan selalu siap membantu.",
             font_size=18, color=LIGHT_GRAY)

# Info pembuat
add_text_box(slide1, Inches(1.2), Inches(5.5), Inches(8), Inches(0.5),
             "Dibuat oleh: Hamzah82  |  GitHub: Hamzah82/Ruka-AI", font_size=14, color=RGBColor(0x88, 0x88, 0x88))

# Tag di kanan bawah
tags = ["Python", "OpenRouter API", "CLI Agent", "Local-First"]
tag_left = Inches(1.2)
for tag in tags:
    tag_shape = add_shape(slide1, tag_left, Inches(6.3), Inches(1.5), Inches(0.4), RGBColor(0x2A, 0x2A, 0x4A), MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tag_shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)
    tf.paragraphs[0].space_after = Pt(2)
    tag_left += Inches(1.7)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — SIAPA RUKA AI?
# ═══════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide2)

# Header
add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Siapa Ruka AI?  🐢", font_size=40, color=WHITE, bold=True)
# Garis bawah header
add_shape(slide2, Inches(0.8), Inches(1.2), Inches(3), Inches(0.05), ACCENT_BLUE, MSO_SHAPE.RECTANGLE)

# Kolom kiri — Definisi
add_text_box(slide2, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
             "Definisi", font_size=24, color=ACCENT_BLUE, bold=True)

definisi_items = [
    "AI Agent berbentuk kura-kura 🐢",
    "Berjalan di terminal (CLI — Command Line Interface)",
    "Terhubung ke OpenRouter API untuk kecerdasan AI",
    "Beroperasi secara lokal di mesin pengguna",
    "Bukan sekadar chatbot — tapi agent yang bisa bertindak",
]
add_bullet_list(slide2, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4), definisi_items, font_size=16)

# Kolom kanan — Karakteristik
add_text_box(slide2, Inches(7), Inches(1.6), Inches(5.5), Inches(0.6),
             "Karakteristik", font_size=24, color=ACCENT_PURPLE, bold=True)

karakter_items = [
    "🤖 Agentic — Bisa memutuskan sendiri tindakan",
    "🏠 Local-first — Semua data di mesin lokal",
    "💾 Session-based — Percakapan tersimpan persisten",
    "🔧 Model-agnostic — Buka pakai model apapun",
    "🛡️ Private — Tidak ada data yang dikirim ke cloud",
]
add_bullet_list(slide2, Inches(7), Inches(2.2), Inches(5.5), Inches(4), karakter_items, font_size=16)

# Quote di bawah
quote_shape = add_shape(slide2, Inches(2), Inches(6.2), Inches(9), Inches(0.6), RGBColor(0x22, 0x22, 0x3A), MSO_SHAPE.ROUNDED_RECTANGLE)
tf = quote_shape.text_frame
p = tf.paragraphs[0]
p.text = '"Bijaksana, sabar, dan teliti — seperti kura-kura yang selalu sampai tujuan."'
p.font.size = Pt(14)
p.font.color.rgb = YELLOW
p.font.italic = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — KEMAMPUAN (TOOLS)
# ═══════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide3)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Kemampuan Ruka AI  🛠️", font_size=40, color=WHITE, bold=True)
add_shape(slide3, Inches(0.8), Inches(1.2), Inches(3), Inches(0.05), ACCENT_GREEN, MSO_SHAPE.RECTANGLE)

# Baris 1 — File Operations
add_text_box(slide3, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
             "📁 Operasi File (Tangan Kanan)", font_size=20, color=ACCENT_ORANGE, bold=True)

file_ops = [
    ("read_file", "Membaca isi file"),
    ("write_file", "Menulis / membuat file"),
    ("delete_file", "Menghapus file"),
    ("copy_file", "Menyalin file"),
    ("move_file", "Memindahkan / rename file"),
    ("edit_file", "Mengedit isi file"),
]

x_start = Inches(0.8)
y_pos = Inches(2.1)
for i, (name, desc) in enumerate(file_ops):
    col = i % 3
    row = i // 3
    x = x_start + col * Inches(4.2)
    y = y_pos + row * Inches(1.1)

    card = add_shape(slide3, x, y, Inches(3.8), Inches(0.85), RGBColor(0x22, 0x22, 0x3A), MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{name}  —  {desc}"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.name = "Consolas"

# Baris 2 — Folder & Info
add_text_box(slide3, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
             "📂 Folder & Informasi (Tangan Kiri)", font_size=20, color=ACCENT_ORANGE, bold=True)

folder_ops = [
    ("create_folder", "Membuat folder baru"),
    ("delete_folder", "Menghapus folder"),
    ("list_all", "Struktur direktori lengkap"),
    ("list_files", "Daftar file"),
    ("get_file_info", "Info detail file / folder"),
    ("exec_command", "Menjalankan perintah terminal"),
]

x_start = Inches(0.8)
y_pos = Inches(5.0)
for i, (name, desc) in enumerate(folder_ops):
    col = i % 3
    row = i // 3
    x = x_start + col * Inches(4.2)
    y = y_pos + row * Inches(1.1)

    card = add_shape(slide3, x, y, Inches(3.8), Inches(0.85), RGBColor(0x22, 0x22, 0x3A), MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{name}  —  {desc}"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.name = "Consolas"


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — ARSITEKTUR SISTEM
# ═══════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide4)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Arsitektur Sistem  🏗️", font_size=40, color=WHITE, bold=True)
add_shape(slide4, Inches(0.8), Inches(1.2), Inches(3), Inches(0.05), ACCENT_PURPLE, MSO_SHAPE.RECTANGLE)

# Komponen-komponen
komponen = [
    ("🧠", "Otak", "Agentic Loop — terhubung ke OpenRouter API untuk AI reasoning"),
    ("🦴", "Tubuh Utama", "Single-file CLI agent (main.py ~87KB, ~2000+ baris)"),
    ("👁️", "Sensor", "Interrupt mechanism — user bisa hentikan proses kapan saja"),
    ("💾", "Memori", "Session system — percakapan tersimpan di sessions/*.json"),
    ("🔒", "Keamanan", "Path traversal protection + command blocklist"),
    ("🌐", "Browsing", "Bisa akses internet via lynx, w3m, curl"),
]

y_pos = Inches(1.7)
for emoji, title, desc in komponen:
    # Card
    card = add_shape(slide4, Inches(0.8), y_pos, Inches(11.5), Inches(0.8), RGBColor(0x22, 0x22, 0x3A), MSO_SHAPE.ROUNDED_RECTANGLE)

    # Emoji + Title
    txBox_title = slide4.shapes.add_textbox(Inches(1.0), y_pos + Inches(0.05), Inches(3), Inches(0.4))
    tf = txBox_title.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji}  {title}"
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_BLUE
    p.font.bold = True
    p.font.name = "Calibri"

    # Description
    txBox_desc = slide4.shapes.add_textbox(Inches(1.0), y_pos + Inches(4), Inches(11), Inches(0.4))
    tf = txBox_desc.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Calibri"

    y_pos += Inches(0.95)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — KEAMANAN
# ═══════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide5)

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Keamanan  🔒", font_size=40, color=WHITE, bold=True)
add_shape(slide5, Inches(0.8), Inches(1.2), Inches(3), Inches(0.05), RGBColor(0xFF, 0x44, 0x44), MSO_SHAPE.RECTANGLE)

# Kolom kiri — Proteksi
add_text_box(slide5, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
             "Jenis Proteksi", font_size=24, color=ACCENT_ORANGE, bold=True)

proteksi_items = [
    "🚫 Path Traversal Protection",
    "   → Semua operasi file dibatasi di dalam BASE_DIR",
    "   → Mencegah akses ke /etc/passwd, ~/.ssh, dll",
    "",
    "🚫 Command Blocklist",
    "   → rm -rf /, mkfs, dd, shutdown, fork bomb",
    "   → Perintah berbahaya diblokir otomatis",
    "",
    "⏱️ Timeout System",
    "   → Setiap perintah terminal punya batas waktu",
    "   → Mencegah infinite loop / hang",
    "",
    "🔑 Environment Variables",
    "   → API key tersimpan di file .env",
    "   → Tidak di-commit ke repository",
]
add_bullet_list(slide5, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), proteksi_items, font_size=14)

# Kolom kanan — Alur keamanan
add_text_box(slide5, Inches(7), Inches(1.6), Inches(5.5), Inches(0.6),
             "Alur Keamanan", font_size=24, color=ACCENT_GREEN, bold=True)

# Flow keamanan
flows = [
    ("1", "User Input", ACCENT_BLUE),
    ("2", "Validasi Path (_safe_path)", ACCENT_PURPLE),
    ("3", "Cek dalam BASE_DIR?", ACCENT_ORANGE),
    ("4a", "Ya → Lanjutkan", ACCENT_GREEN),
    ("4b", "Tidak → Tolak Akses", RGBColor(0xFF, 0x44, 0x44)),
    ("5", "Eksekusi Tool", ACCENT_BLUE),
    ("6", "Return Hasil", ACCENT_GREEN),
]

y_flow = Inches(2.2)
for num, desc, color in flows:
    circle = add_circle(slide5, Inches(7.2), y_flow + Inches(0.05), Inches(0.35), color)
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide5, Inches(7.7), y_flow, Inches(4.5), Inches(0.45),
                 desc, font_size=14, color=color)
    y_flow += Inches(0.55)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — AGENTIC LOOP (CARA KERJA)
# ═══════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide6)

add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Agentic Loop  — Cara Kerja  🔄", font_size=36, color=WHITE, bold=True)
add_shape(slide6, Inches(0.8), Inches(1.2), Inches(3), Inches(0.05), ACCENT_BLUE, MSO_SHAPE.RECTANGLE)

# Flow diagram
steps = [
    ("User Input", "User mengetik perintah", ACCENT_BLUE),
    ("System Prompt + History", "Dikirim ke OpenRouter API", ACCENT_PURPLE),
    ("AI Response", "API mengembalikan response", ACCENT_ORANGE),
    ("Text Response?", "Jika ya → tampilkan ke user", ACCENT_GREEN),
    ("Tool Calls?", "Jika ada → eksekusi semua tool", ACCENT_ORANGE),
    ("Loop", "Hasil tool dikirim lagi ke API", ACCENT_BLUE),
]

y_step = Inches(1.6)
for i, (title, desc, color) in enumerate(steps):
    # Kotak
    box = add_shape(slide6, Inches(2), y_step, Inches(9), Inches(0.75), RGBColor(0x22, 0x22, 0x3A), MSO_SHAPE.ROUNDED_RECTANGLE)

    # Border kiri berwarna
    left_bar = add_shape(slide6, Inches(2), y_step, Inches(0.08), Inches(0.75), color, MSO_SHAPE.RECTANGLE)

    # Title
    txBox = slide6.shapes.add_textbox(Inches(2.3), y_step + Inches(0.05), Inches(8), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"[{i+1}] {title}"
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"

    # Desc
    txBox2 = slide6.shapes.add_textbox(Inches(2.3), y_step + Inches(3.5), Inches(8), Inches(0.35))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Calibri"

    # Arrow antar step
    if i < len(steps) - 1:
        arrow = add_shape(slide6, Inches(6.3), y_step + Inches(0.75), Inches(0.5), Inches(0.15), RGBColor(0x44, 0x44, 0x6A), MSO_SHAPE.DOWN_ARROW)
        tf_arrow = arrow.text_frame
        p_arrow = tf_arrow.paragraphs[0]
        p_arrow.text = "▼"
        p_arrow.font.size = Pt(14)
        p_arrow.font.color.rgb = RGBColor(0x66, 0x66, 0x88)
        p_arrow.alignment = PP_ALIGN.CENTER

    y_step += Inches(0.9)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — FITUR UNGGULAN
# ═══════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide7)

add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Fitur Unggulan  ⭐", font_size=40, color=WHITE, bold=True)
add_shape(slide7, Inches(0.8), Inches(1.2), Inches(3), Inches(0.05), YELLOW, MSO_SHAPE.RECTANGLE)

fitur = [
    ("🐢", "Multi-Step Execution", "Bisa pecah tugas kompleks menjadi beberapa langkah berturut-turut"),
    ("🔧", "12 Tools Lengkap", "File, folder, info, terminal — semua dalam satu agent"),
    ("💾", "Session Persisten", "Percakapan tersimpan otomatis, bisa dilanjutkan nanti"),
    ("🌐", "Internet Access", "Bisa browsing, search, dan scrape web"),
    ("🛡️", "Aman & Private", "Semua berjalan lokal, data tidak dikirim ke cloud"),
    ("🔄", "Interruptible", "User bisa hentikan proses kapan saja dengan ketik 'q'"),
    ("📝", "Smart File Edit", "Replace, append, presisi — edit file tanpa tulis ulang"),
    ("🎨", "No Dependencies", "Hanya butuh Python + 1 package (python-pptx untuk ekstra)"),
]

# 2 kolom x 4 baris
for i, (emoji, title, desc) in enumerate(fitur):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6)
    y = Inches(1.7) + row * Inches(1.35)

    card = add_shape(slide7, x, y, Inches(5.5), Inches(1.15), RGBColor(0x22, 0x22, 0x3A), MSO_SHAPE.ROUNDED_RECTANGLE)

    # Emoji + Title
    txBox_t = slide7.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), Inches(5), Inches(0.4))
    tf = txBox_t.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji}  {title}"
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_BLUE
    p.font.bold = True
    p.font.name = "Calibri"

    # Desc
    txBox_d = slide7.shapes.add_textbox(x + Inches(0.2), y + Inches(0.55), Inches(5), Inches(0.5))
    tf2 = txBox_d.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — TEKNOLOGI
# ═══════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide8)

add_text_box(slide8, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Teknologi yang Digunakan  💻", font_size=40, color=WHITE, bold=True)
add_shape(slide8, Inches(0.8), Inches(1.2), Inches(3), Inches(0.05), ACCENT_GREEN, MSO_SHAPE.RECTANGLE)

# Stack teknologi
tech_stack = [
    ("🐍", "Python 3.13", "Bahasa pemrograman utama — single file architecture"),
    ("🤖", "OpenRouter API", "Gateway ke berbagai model AI (GPT, Claude, Gemini, dll)"),
    ("📦", "python-pptx", "Library untuk generate file PPT"),
    ("🖼️", "Pillow", "Image processing support"),
    ("🌐", "lynx / w3m / curl", "Text-based web browsing & scraping"),
    ("📋", "JSON", "Format penyimpanan session & data"),
    ("🔑", ".env", "Environment variables untuk API keys"),
    ("📁", "shutil / os", "File & folder operations"),
]

y_tech = Inches(1.7)
for emoji, name, desc in tech_stack:
    card = add_shape(slide8, Inches(1.5), y_tech, Inches(10), Inches(0.65), RGBColor(0x22, 0x22, 0x3A), MSO_SHAPE.ROUNDED_RECTANGLE)

    txBox = slide8.shapes.add_textbox(Inches(1.8), y_tech + Inches(0.05), Inches(9.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji}  {name}"
    p.font.size = Pt(16)
    p.font.color.rgb = ACCENT_GREEN
    p.font.bold = True
    p.font.name = "Calibri"

    txBox2 = slide8.shapes.add_textbox(Inches(1.8), y_tech + Inches(0.32), Inches(9.5), Inches(0.3))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Calibri"

    y_tech += Inches(0.75)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — PENUTUP
# ═══════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide9)

# Dekorasi
add_circle(slide9, Inches(8), Inches(-2), Inches(6), RGBColor(0x16, 0x21, 0x37))
add_circle(slide9, Inches(-1), Inches(4), Inches(5), RGBColor(0x16, 0x21, 0x37))

# Judul
add_text_box(slide9, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1),
             "Terima Kasih!  🐢", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Sub
add_text_box(slide9, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.8),
             "Ruka AI — Kura-Kura Pintar di Terminal Kamu", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Links
add_text_box(slide9, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5),
             "📂 GitHub:  github.com/Hamzah82/Ruka-AI", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Quote
quote2 = add_shape(slide9, Inches(2.5), Inches(4.8), Inches(8), Inches(0.8), RGBColor(0x22, 0x22, 0x3A), MSO_SHAPE.ROUNDED_RECTANGLE)
tf = quote2.text_frame
p = tf.paragraphs[0]
p.text = '"Pelan tapi pasti — seperti kura-kura, Ruka AI selalu sampai tujuan."'
p.font.size = Pt(16)
p.font.color.rgb = YELLOW
p.font.italic = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Footer
add_text_box(slide9, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
             "Dibuat dengan ❤️ oleh Hamzah82  |  © 2025 Ruka AI", font_size=12, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SIMPAN FILE
# ═══════════════════════════════════════════════════════════
output_path = "Ruka_AI_Pengenalan.pptx"
prs.save(output_path)
print(f"✅ File PPT berhasil dibuat: {output_path}")
print(f"📊 Jumlah slide: {len(prs.slides)}")
print(f"📐 Ukuran: {SLIDE_W}, {SLIDE_H}")
